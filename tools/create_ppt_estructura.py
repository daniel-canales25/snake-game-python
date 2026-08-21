"""
Script para generar presentación PowerPoint: Estructura y Retos del Snake Game.
Uso: uv run python tools/create_ppt_estructura.py
"""
import subprocess
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = "tools/temp_images"
OUTPUT_FILE = "snake_game_estructura.pptx"

# Colores del tema
BG_COLOR = RGBColor(25, 25, 35)
TITLE_COLOR = RGBColor(0, 200, 83)
TEXT_COLOR = RGBColor(255, 255, 255)
ACCENT_COLOR = RGBColor(100, 180, 255)
CHALLENGE_COLOR = RGBColor(255, 180, 50)
SUBTITLE_COLOR = RGBColor(180, 180, 180)

# Diagramas Mermaid
DIAGRAMS = {
    "project_structure": """
flowchart TD
    subgraph Proyecto
        A[snake-game-python] --> B[src/]
        A --> C[tests/]
        A --> D[tools/]
        A --> E[codigo-mermaid/]
        B --> F[snake_game/]
        F --> G[core/]
        F --> H[utils/]
        G --> I[game.py]
        G --> J[snake.py]
        G --> K[food.py]
        G --> L[start_screen.py]
        G --> M[name_input.py]
        H --> N[constants.py]
        H --> O[score_manager.py]
    end
    """,
    "game_flow": """
flowchart TD
    A([Inicio]) --> B[Pygame Init]
    B --> C[Pantalla Inicio]
    C --> D[Leaderboard + JUGAR]
    D --> E{Click JUGAR}
    E --> F[Juego Activo]
    F --> G[Mover serpiente]
    G --> H{Comer comida?}
    H -->|Yes| I[Crecer + Score]
    H -->|No| J{Colisión?}
    I --> G
    J -->|No| G
    J -->|Yes| K[Game Over]
    K --> L[Ingresar Nombre]
    L --> M[Guardar Score]
    M --> C
    """,
    "snake_class": """
flowchart TD
    A[SNAKE] --> B[body: lista de tuplas]
    A --> C[direction: tupla x,y]
    A --> D[nextDirection: buffer]
    A --> E[isGrowing: boolean]
    B --> F[head = body 0]
    A --> G[move]
    A --> H[grow]
    A --> I[change_direction]
    A --> J[draw]
    G --> K[Insertar nueva cabeza]
    G --> L[Eliminar cola si no crece]
    H --> M[isGrowing = True]
    I --> N[Rechazar dirección opuesta]
    J --> O[Dibujar rectángulos]
    """,
    "food_class": """
flowchart TD
    A[FOOD] --> B[position: tupla x,y]
    A --> C[respawn]
    A --> D[draw]
    C --> E[Generar posición aleatoria]
    E --> F{En snakeBody?}
    F -->|Yes| E
    F -->|No| G[Asignar position]
    D --> H[Dibujar círculo rojo]
    """,
    "game_class": """
flowchart TD
    A[GAME] --> B[snake: Snake]
    A --> C[food: Food]
    A --> D[score: int]
    A --> E[isGameOver: bool]
    A --> F[update]
    A --> G[draw]
    F --> H[Mover snake]
    H --> I{Comida comida?}
    I -->|Si| J[Crecer + Score + Respawn]
    I -->|No| K{Colision pared?}
    K -->|Si| L[Game Over]
    K -->|No| M[Continuar]
    J --> M
    """
}


def setup_temp_dir():
    os.makedirs(OUTPUT_DIR, exist_ok=True)


def mermaid_to_image(mermaid_code, output_name):
    temp_mmd = os.path.join(OUTPUT_DIR, f"{output_name}.mmd")
    temp_png = os.path.join(OUTPUT_DIR, f"{output_name}.png")

    with open(temp_mmd, 'w', encoding='utf-8') as f:
        f.write(mermaid_code)

    cmd = [
        "npx", "--yes", "@mermaid-js/mermaid-cli",
        "-i", temp_mmd,
        "-o", temp_png,
        "-b", "transparent",
        "-w", "1200"
    ]

    try:
        subprocess.run(cmd, capture_output=True, check=True, timeout=30)
        print(f"  ✓ {temp_png}")
        return temp_png
    except Exception as e:
        print(f"  ✗ Error: {e}")
        return None


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, top=Inches(0.4), font_size=Pt(36)):
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = TITLE_COLOR
    p.font.bold = True
    return txBox


def add_text(slide, text, left, top, width, height, font_size=Pt(18), color=TEXT_COLOR, alignment=PP_ALIGN.LEFT, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, items, left, top, width, height, font_size=Pt(16), color=TEXT_COLOR):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = font_size
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return txBox


def add_image(slide, image_path, left, top, width=None):
    if width:
        slide.shapes.add_picture(image_path, left, top, width=width)
    else:
        slide.shapes.add_picture(image_path, left, top)


def add_numbered_item(slide, number, title, description, left, top, width):
    """Agrega un elemento numerado con título y descripción."""
    # Número
    add_text(slide, str(number), left, top, Inches(0.6), Inches(0.5),
             font_size=Pt(28), color=TITLE_COLOR, bold=True)
    # Título
    add_text(slide, title, left + Inches(0.7), top, width - Inches(0.7), Inches(0.4),
             font_size=Pt(18), color=ACCENT_COLOR, bold=True)
    # Descripción
    add_text(slide, description, left + Inches(0.7), top + Inches(0.4), width - Inches(0.7), Inches(0.5),
             font_size=Pt(14), color=SUBTITLE_COLOR)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # === DIAPOSITIVA 1: PORTADA ===
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide1, BG_COLOR)

    add_title(slide1, "SNAKE GAME", top=Inches(2), font_size=Pt(54))
    add_text(slide1, "Estructura del Proyecto y Retos",
             Inches(0.5), Inches(3.2), Inches(9), Inches(0.8),
             font_size=Pt(24), color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)
    add_text(slide1, "Lógica de Programación",
             Inches(0.5), Inches(4), Inches(9), Inches(0.8),
             font_size=Pt(18), color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)

    # === DIAPOSITIVA 2: ESTRUCTURA DEL PROYECTO ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide2, BG_COLOR)

    add_title(slide2, "Estructura del Proyecto")

    if "project_structure" in DIAGRAMS:
        img_path = mermaid_to_image(DIAGRAMS["project_structure"], "project_structure")
        if img_path:
            add_image(slide2, img_path, Inches(0.5), Inches(1.5), width=Inches(9))

    items = [
        "src/snake_game/ → Código fuente principal",
        "core/ → Clases del juego (Snake, Food, Game, etc.)",
        "utils/ → Constantes y gestor de scores",
        "tools/ → Scripts auxiliares (generador de presentación)",
        "codigo-mermaid/ → Diagramas de flujo"
    ]
    add_bullet_list(slide2, items, Inches(0.5), Inches(5.5), Inches(9), Inches(2))

    # === DIAPOSITIVA 3: MÁQUINA DE ESTADOS ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide3, BG_COLOR)

    add_title(slide3, "Flujo del Juego (Máquina de Estados)")

    if "game_flow" in DIAGRAMS:
        img_path = mermaid_to_image(DIAGRAMS["game_flow"], "game_flow")
        if img_path:
            add_image(slide3, img_path, Inches(0.5), Inches(1.5), width=Inches(9))

    items = [
        "3 estados: START_SCREEN → PLAYING → NAME_INPUT",
        "El juego comienza en Pantalla de Inicio",
        "Al ganar, se guarda score y se vuelve al inicio"
    ]
    add_bullet_list(slide3, items, Inches(0.5), Inches(5.8), Inches(9), Inches(1.5))

    # === DIAPOSITIVA 4: CLASE SNAKE ===
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide4, BG_COLOR)

    add_title(slide4, "Clase: Snake")

    if "snake_class" in DIAGRAMS:
        img_path = mermaid_to_image(DIAGRAMS["snake_class"], "snake_class")
        if img_path:
            add_image(slide4, img_path, Inches(0.3), Inches(1.5), width=Inches(5))

    items = [
        "body: Lista de tuplas (x, y)",
        "direction: Tupla (1,0) derecha",
        "move(): Inserta cabeza, elimina cola",
        "grow(): Activa crecimiento",
        "draw(): Rectángulos verdes en grid"
    ]
    add_bullet_list(slide4, items, Inches(5.5), Inches(1.5), Inches(4.2), Inches(3))

    add_text(slide4, "Grid: 30x30 celdas de 33px",
             Inches(0.5), Inches(6.5), Inches(9), Inches(0.5),
             font_size=Pt(14), color=ACCENT_COLOR)

    # === DIAPOSITIVA 5: CLASES FOOD Y GAME ===
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide5, BG_COLOR)

    add_title(slide5, "Clases: Food y Game")

    # Food
    if "food_class" in DIAGRAMS:
        img_path = mermaid_to_image(DIAGRAMS["food_class"], "food_class")
        if img_path:
            add_image(slide5, img_path, Inches(0.2), Inches(1.5), width=Inches(4.5))

    add_text(slide5, "FOOD",
             Inches(0.5), Inches(1.2), Inches(4), Inches(0.4),
             font_size=Pt(18), color=TITLE_COLOR, bold=True)

    # Game
    if "game_class" in DIAGRAMS:
        img_path = mermaid_to_image(DIAGRAMS["game_class"], "game_class")
        if img_path:
            add_image(slide5, img_path, Inches(5.2), Inches(1.5), width=Inches(4.5))

    add_text(slide5, "GAME",
             Inches(5.5), Inches(1.2), Inches(4), Inches(0.4),
             font_size=Pt(18), color=TITLE_COLOR, bold=True)

    items_food = [
        "Posición aleatoria en grid",
        "Respawn sin colisión con snake",
        "Círculo rojo de 33px"
    ]
    add_bullet_list(slide5, items_food, Inches(0.5), Inches(4.5), Inches(4.2), Inches(2))

    items_game = [
        "Orquesta snake + food",
        "Detecta colisiones",
        "Maneja score y game over"
    ]
    add_bullet_list(slide5, items_game, Inches(5.5), Inches(4.5), Inches(4.2), Inches(2))

    # === DIAPOSITIVA 6: CLASES UI ===
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide6, BG_COLOR)

    add_title(slide6, "Clases: StartScreen y NameInput")

    # StartScreen
    add_text(slide6, "STARTSCREEN",
             Inches(0.5), Inches(1.2), Inches(4.5), Inches(0.4),
             font_size=Pt(18), color=TITLE_COLOR, bold=True)

    items_start = [
        "Muestra leaderboard (top 10 scores)",
        "Botón JUGAR con efecto hover",
        "Mensaje si no calificó al top 10",
        "Scores persistidos en scores.json"
    ]
    add_bullet_list(slide6, items_start, Inches(0.5), Inches(1.8), Inches(4.2), Inches(2.5))

    # NameInput
    add_text(slide6, "NAMEINPUT",
             Inches(5.5), Inches(1.2), Inches(4.2), Inches(0.4),
             font_size=Pt(18), color=TITLE_COLOR, bold=True)

    items_name = [
        "Pantalla post game over",
        "Input de nombre (max 5 chars)",
        "Cursor parpadeante",
        "ENTER guarda / SPACE reinicia"
    ]
    add_bullet_list(slide6, items_name, Inches(5.5), Inches(1.8), Inches(4.2), Inches(2.5))

    add_text(slide6, "Constantes: 30x30 grid, 7 FPS, colores definidos en constants.py",
             Inches(0.5), Inches(6.5), Inches(9), Inches(0.5),
             font_size=Pt(14), color=ACCENT_COLOR)

    # === DIAPOSITIVA 7: RETOS ===
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide7, BG_COLOR)

    add_title(slide7, "Retos del Proyecto")

    challenges = [
        (1, "UV Python",
         "Aprender a usar uv como gestor de dependencias moderno. Configurar el entorno virtual, agregar/quitar paquetes con uv add/remove."),
        (2, "Configuración de Pygame",
         "Entender la inicialización de Pygame, manejo de eventos, loop del juego a 7 FPS, y dibujado en pantalla con pygame.draw."),
        (3, "Herramientas de Diagramas",
         "Aprender Draw.io para crear diagramas de flujo y Mermaid para documentar la lógica del código de forma visual."),
        (4, "Primer Proyecto Python",
         "Mi primer proyecto completo en Python: estructura de carpetas, módulos, clases, y persistencia de datos con JSON.")
    ]

    y_pos = Inches(1.6)
    for num, title, desc in challenges:
        add_numbered_item(slide7, num, title, desc, Inches(0.5), y_pos, Inches(9))
        y_pos += Inches(1.3)

    # === DIAPOSITIVA 8: CONCLUSIÓN ===
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide8, BG_COLOR)

    add_title(slide8, "Conclusión")

    items = [
        "Estructura clara: core/ para lógica, utils/ para helpers",
        "Máquina de estados controla el flujo del juego",
        "Persistencia de scores con JSON",
        "Documentación visual con Mermaid",
        "Herramientas modernas: uv, Pygame, Draw.io"
    ]
    add_bullet_list(slide8, items, Inches(0.5), Inches(1.8), Inches(9), Inches(3))

    add_text(slide8, "¡Proyecto completado con éxito!",
             Inches(0.5), Inches(5.5), Inches(9), Inches(0.8),
             font_size=Pt(24), color=TITLE_COLOR, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide8, "snake-game-python/",
             Inches(0.5), Inches(6.3), Inches(9), Inches(0.5),
             font_size=Pt(16), color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)

    # Guardar
    prs.save(OUTPUT_FILE)
    print(f"\n✓ Presentación guardada: {OUTPUT_FILE}")


if __name__ == "__main__":
    print("=== Generando presentación: Estructura y Retos ===\n")
    setup_temp_dir()
    create_presentation()

    import shutil
    if os.path.exists(OUTPUT_DIR):
        shutil.rmtree(OUTPUT_DIR)
        print("✓ Archivos temporales eliminados")
