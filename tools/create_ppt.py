"""
Script para generar presentación PowerPoint del Snake Game.
Uso: uv run python tools/create_ppt.py
"""
import subprocess
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

MERMAID_DIR = "codigo-mermaid"
OUTPUT_DIR = "tools/temp_images"
OUTPUT_FILE = "snake_game_presentation.pptx"

# Colores del tema
BG_COLOR = RGBColor(30, 30, 30)
TITLE_COLOR = RGBColor(0, 200, 83)
TEXT_COLOR = RGBColor(255, 255, 255)
ACCENT_COLOR = RGBColor(100, 180, 255)

# Diagramas Mermaid simplificados para la presentación
DIAGRAMS = {
    "architecture": """
flowchart TD
    subgraph Estructura del Proyecto
        A[src/snake_game] --> B[core/]
        A --> C[utils/]
        B --> D[game.py]
        B --> E[snake.py]
        B --> F[food.py]
        B --> G[start_screen.py]
        B --> H[name_input.py]
        C --> I[constants.py]
        C --> J[score_manager.py]
    end
    K[main.py] --> A
    """,
    "state_machine": """
flowchart LR
    A[START_SCREEN] -->|Click JUGAR o SPACE| B[PLAYING]
    B -->|Game Over| C[NAME_INPUT]
    C -->|ENTER guarda score| A
    C -->|SPACE reinicia| B
    """,
    "game_loop": """
flowchart TD
    A([Inicio]) --> B[Pygame Init]
    B --> C[Crear screen + clock]
    C --> D{running?}
    D -->|Yes| E[Procesar eventos]
    E --> F[Dibujar estado actual]
    F --> G[pygame.display.flip]
    G --> H[clock.tick 7 FPS]
    H --> D
    D -->|No| I[pygame.quit]
    I --> J([Fin])
    """,
    "snake_move": """
flowchart TD
    A([move]) --> B[Calcular nueva cabeza]
    B --> C{Y < scoreAreaHeight?}
    C -->|Yes| D[Clamp Y = 3]
    C -->|No| E[Insertar cabeza]
    D --> E
    E --> F{isGrowing?}
    F -->|No| G[Eliminar cola]
    F -->|Yes| H[isGrowing = False]
    G --> I([Fin])
    H --> I
    """,
    "food_respawn": """
flowchart TD
    A([respawn]) --> B[Generar posición aleatoria]
    B --> C{Posición en snakeBody?}
    C -->|Yes| B
    C -->|No| D[Asignar position]
    D --> E([Fin])
    """,
    "draw_overview": """
flowchart TD
    A[draw] --> B{current_state?}
    B -->|START_SCREEN| C[start_screen.draw]
    B -->|PLAYING| D[game.update + draw]
    B -->|NAME_INPUT| E[name_input.draw]
    C --> F([Fin])
    D --> F
    E --> F
    """
}


def setup_temp_dir():
    """Crea directorio temporal para imágenes."""
    os.makedirs(OUTPUT_DIR, exist_ok=True)


def extract_mermaid_code(md_file):
    """Extrae código Mermaid de un archivo .md."""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    start = content.find('```mermaid') + 9
    end = content.find('```', start)
    return content[start:end].strip()


def mermaid_to_image(mermaid_code, output_name):
    """Convierte código Mermaid a imagen PNG usando mermaid-cli."""
    temp_mmd = os.path.join(OUTPUT_DIR, f"{output_name}.mmd")
    temp_png = os.path.join(OUTPUT_DIR, f"{output_name}.png")

    with open(temp_mmd, 'w', encoding='utf-8') as f:
        f.write(mermaid_code)

    cmd = [
        "npx", "--yes", "@mermaid-js/mermaid-cli",
        "-i", temp_mmd,
        "-o", temp_png,
        "-b", "transparent",
        "-w", "1200"
    ]

    try:
        subprocess.run(cmd, capture_output=True, check=True, timeout=30)
        print(f"  ✓ Generada: {temp_png}")
        return temp_png
    except subprocess.TimeoutExpired:
        print(f"  ✗ Timeout generando: {output_name}")
        return None
    except subprocess.CalledProcessError as e:
        print(f"  ✗ Error generando {output_name}: {e}")
        return None


def set_slide_bg(slide, color):
    """Establece el fondo de una diapositiva."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, top=Inches(0.5), font_size=Pt(36)):
    """Agrega un título a la diapositiva."""
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = TITLE_COLOR
    p.font.bold = True
    return txBox


def add_text(slide, text, left, top, width, height, font_size=Pt(18), color=TEXT_COLOR, alignment=PP_ALIGN.LEFT):
    """Agrega texto a la diapositiva."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, items, left, top, width, height, font_size=Pt(16)):
    """Agrega una lista de viñetas."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = font_size
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)
    return txBox


def add_image(slide, image_path, left, top, width=None, height=None):
    """Agrega una imagen a la diapositiva."""
    if width and height:
        slide.shapes.add_picture(image_path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(image_path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(image_path, left, top, height=height)
    else:
        slide.shapes.add_picture(image_path, left, top)


def create_presentation():
    """Crea la presentación PowerPoint."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # === DIAPOSITIVA 1: PORTADA ===
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide1, BG_COLOR)

    add_title(slide1, "SNAKE GAME", top=Inches(2), font_size=Pt(54))
    add_text(slide1, "Python + Pygame",
             Inches(0.5), Inches(3.2), Inches(9), Inches(0.8),
             font_size=Pt(28), color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)
    add_text(slide1, "Lógica de Programación",
             Inches(0.5), Inches(4), Inches(9), Inches(0.8),
             font_size=Pt(20), color=TEXT_COLOR, alignment=PP_ALIGN.CENTER)

    # === DIAPOSITIVA 2: ARQUITECTURA ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide2, BG_COLOR)

    add_title(slide2, "Arquitectura del Proyecto")

    if "architecture" in DIAGRAMS:
        img_path = mermaid_to_image(DIAGRAMS["architecture"], "architecture")
        if img_path:
            add_image(slide2, img_path, Inches(1), Inches(1.8), width=Inches(8))

    add_text(slide2, "src/snake_game/",
             Inches(0.5), Inches(6), Inches(9), Inches(0.5),
             font_size=Pt(14), color=ACCENT_COLOR)

    # === DIAPOSITIVA 3: MÁQUINA DE ESTADOS ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide3, BG_COLOR)

    add_title(slide3, "Máquina de Estados")

    if "state_machine" in DIAGRAMS:
        img_path = mermaid_to_image(DIAGRAMS["state_machine"], "state_machine")
        if img_path:
            add_image(slide3, img_path, Inches(0.5), Inches(1.8), width=Inches(9))

    items = [
        "START_SCREEN: Pantalla de inicio con leaderboard",
        "PLAYING: Juego activo (movimiento + colisiones)",
        "NAME_INPUT: Ingreso de nombre post game over"
    ]
    add_bullet_list(slide3, items, Inches(0.5), Inches(5), Inches(9), Inches(2))

    # === DIAPOSITIVA 4: GAME LOOP ===
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide4, BG_COLOR)

    add_title(slide4, "Flujo del Juego (Game Loop)")

    if "game_loop" in DIAGRAMS:
        img_path = mermaid_to_image(DIAGRAMS["game_loop"], "game_loop")
        if img_path:
            add_image(slide4, img_path, Inches(1.5), Inches(1.8), width=Inches(7))

    items = [
        "Ejecuta a 7 FPS (7 iteraciones por segundo)",
        "Procesa eventos de teclado y mouse",
        "Dibuja el estado actual en pantalla"
    ]
    add_bullet_list(slide4, items, Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))

    # === DIAPOSITIVA 5: CLASES PRINCIPALES ===
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide5, BG_COLOR)

    add_title(slide5, "Clases Principales")

    classes = [
        ("Snake", "Movimiento, crecimiento, dirección, dibujado"),
        ("Food", "Posicionamiento aleatorio, respawn sin colisión"),
        ("Game", "Orquestador: colisiones, score, update/draw"),
        ("StartScreen", "Leaderboard + botón JUGAR"),
        ("NameInput", "Entrada de nombre (max 5 chars)")
    ]

    y_pos = Inches(1.8)
    for class_name, desc in classes:
        add_text(slide5, class_name, Inches(0.5), y_pos, Inches(2.5), Inches(0.5),
                 font_size=Pt(20), color=TITLE_COLOR)
        add_text(slide5, desc, Inches(3), y_pos, Inches(6.5), Inches(0.5),
                 font_size=Pt(16), color=TEXT_COLOR)
        y_pos += Inches(0.7)

    # === DIAPOSITIVA 6: SNAKE Y FOOD ===
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide6, BG_COLOR)

    add_title(slide6, "Snake y Food")

    if "snake_move" in DIAGRAMS:
        img_path = mermaid_to_image(DIAGRAMS["snake_move"], "snake_move")
        if img_path:
            add_image(slide6, img_path, Inches(0.2), Inches(1.8), width=Inches(4.5))

    if "food_respawn" in DIAGRAMS:
        img_path = mermaid_to_image(DIAGRAMS["food_respawn"], "food_respawn")
        if img_path:
            add_image(slide6, img_path, Inches(5.2), Inches(1.8), width=Inches(4.5))

    add_text(slide6, "Snake.move()",
             Inches(0.5), Inches(1.5), Inches(4), Inches(0.4),
             font_size=Pt(14), color=ACCENT_COLOR)
    add_text(slide6, "Food.respawn()",
             Inches(5.5), Inches(1.5), Inches(4), Inches(0.4),
             font_size=Pt(14), color=ACCENT_COLOR)

    # === DIAPOSITIVA 7: DIBUJADO Y CONCLUSIÓN ===
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide7, BG_COLOR)

    add_title(slide7, "Conclusión")

    items = [
        "Máquina de 3 estados controla el flujo",
        "Snake se mueve en grid 30x30 (celdas de 33px)",
        "Colisiones: pared,自身, comida",
        "Score se persiste en scores.json (top 10)",
        "Loop a 7 FPS con Pygame"
    ]
    add_bullet_list(slide7, items, Inches(0.5), Inches(1.8), Inches(9), Inches(3))

    add_text(slide7, "Código fuente: src/snake_game/",
             Inches(0.5), Inches(5.5), Inches(9), Inches(0.5),
             font_size=Pt(16), color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)

    # Guardar presentación
    prs.save(OUTPUT_FILE)
    print(f"\n✓ Presentación guardada: {OUTPUT_FILE}")


if __name__ == "__main__":
    print("=== Generando presentación Snake Game ===\n")
    setup_temp_dir()
    create_presentation()

    # Limpiar archivos temporales
    import shutil
    if os.path.exists(OUTPUT_DIR):
        shutil.rmtree(OUTPUT_DIR)
        print("✓ Archivos temporales eliminados")
