"""
Script para generar presentación PowerPoint: Solo texto con puntos clave.
Uso: uv run python tools/create_ppt_texto.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_FILE = "snake_game_puntos_clave.pptx"

# Colores
BG_COLOR = RGBColor(25, 25, 35)
TITLE_COLOR = RGBColor(0, 200, 83)
TEXT_COLOR = RGBColor(255, 255, 255)
ACCENT_COLOR = RGBColor(100, 180, 255)
SUBTITLE_COLOR = RGBColor(180, 180, 180)


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, top=Inches(0.5), font_size=Pt(36)):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = TITLE_COLOR
    p.font.bold = True


def add_subtitle(slide, text, top=Inches(1.2), font_size=Pt(20)):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = ACCENT_COLOR


def add_text(slide, text, left, top, width, height, font_size=Pt(18), color=TEXT_COLOR, alignment=PP_ALIGN.LEFT, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment


def add_bullet_list(slide, items, left, top, width, height, font_size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = font_size
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(12)


def add_two_columns(slide, left_title, left_items, right_title, right_items, top=Inches(2)):
    # Columna izquierda
    add_text(slide, left_title, Inches(0.5), top, Inches(4), Inches(0.5),
             font_size=Pt(20), color=ACCENT_COLOR, bold=True)
    add_bullet_list(slide, left_items, Inches(0.5), top + Inches(0.5), Inches(4), Inches(3),
                    font_size=Pt(16))

    # Columna derecha
    add_text(slide, right_title, Inches(5.5), top, Inches(4), Inches(0.5),
             font_size=Pt(20), color=ACCENT_COLOR, bold=True)
    add_bullet_list(slide, right_items, Inches(5.5), top + Inches(0.5), Inches(4), Inches(3),
                    font_size=Pt(16))


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # === DIAPOSITIVA 1: PORTADA ===
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide1, BG_COLOR)
    add_title(slide1, "SNAKE GAME", top=Inches(2.5), font_size=Pt(54))
    add_text(slide1, "Python + Pygame", Inches(0.5), Inches(3.5), Inches(9), Inches(0.8),
             font_size=Pt(28), color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)
    add_text(slide1, "Lógica de Programación", Inches(0.5), Inches(4.3), Inches(9), Inches(0.8),
             font_size=Pt(20), color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)

    # === DIAPOSITIVA 2: ¿QUÉ ES EL JUEGO? ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide2, BG_COLOR)
    add_title(slide2, "¿Qué es el Snake Game?")
    add_bullet_list(slide2, [
        "Juego clásico donde controlas una serpiente",
        "La serpiente se mueve en una cuadrícula 30x30",
        "El objetivo es comer comida para crecer",
        "El juego termina si chocas con una pared o contigo mismo",
        "Los scores se guardan en un leaderboard top 10"
    ], Inches(0.5), Inches(1.8), Inches(9), Inches(4))

    # === DIAPOSITIVA 3: ESTRUCTURA DEL PROYECTO ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide3, BG_COLOR)
    add_title(slide3, "Estructura del Proyecto")

    add_two_columns(slide3,
        "Código Fuente", [
            "src/snake_game/",
            "core/ → Clases del juego",
            "utils/ → Constantes y helpers",
            "main.py → Punto de entrada"
        ],
        "Archivos Importantes", [
            "constants.py → Configuración",
            "score_manager.py → Persistencia",
            "scores.json → Datos guardados",
            "pyproject.toml → Dependencias"
        ]
    )

    # === DIAPOSITIVA 4: MÁQUINA DE ESTADOS ===
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide4, BG_COLOR)
    add_title(slide4, "Máquina de Estados")
    add_subtitle(slide4, "El juego tiene 3 estados principales")

    add_bullet_list(slide4, [
        "START_SCREEN → Pantalla de inicio con leaderboard",
        "PLAYING → Juego activo (movimiento y colisiones)",
        "NAME_INPUT → Ingreso de nombre después del game over",
        "",
        "Flujo: START_SCREEN → PLAYING → NAME_INPUT → START_SCREEN",
        "El estado controla qué se dibuja y qué eventos se procesan"
    ], Inches(0.5), Inches(2), Inches(9), Inches(4.5))

    # === DIAPOSITIVA 5: CLASES PRINCIPALES ===
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide5, BG_COLOR)
    add_title(slide5, "Clases Principales")

    add_two_columns(slide5,
        "Snake y Food", [
            "Snake: body, direction, move(), grow()",
            "Food: position, respawn(), draw()",
            "Snake se mueve en grid de 33px por celda",
            "Food aparece en posición aleatoria"
        ],
        "Game y UI", [
            "Game: orquesta snake + food + colisiones",
            "StartScreen: leaderboard + botón JUGAR",
            "NameInput: input de nombre (max 5 chars)",
            "Score se incrementa +10 por comida"
        ]
    )

    # === DIAPOSITIVA 6: FLUJO DEL JUEGO ===
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide6, BG_COLOR)
    add_title(slide6, "Flujo del Juego")

    add_bullet_list(slide6, [
        "1. Pygame init + crear ventana 990x990",
        "2. Mostrar pantalla de inicio con leaderboard",
        "3. Jugador presiona JUGAR o ESPACIO",
        "4. Crear snake en posición inicial (5,8)",
        "5. Snake se mueve a 7 FPS (7 ticks por segundo)",
        "6. Detectar colisiones: comida, pared,自身",
        "7. Si come: crecer +10 score, respawn food",
        "8. Si muere: mostrar input de nombre",
        "9. Guardar score y volver al inicio"
    ], Inches(0.5), Inches(1.8), Inches(9), Inches(5))

    # === DIAPOSITIVA 7: DETALLES TÉCNICOS ===
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide7, BG_COLOR)
    add_title(slide7, "Detalles Técnicos")

    add_two_columns(slide7,
        "Configuración", [
            "Ventana: 990x990 píxeles",
            "Grid: 30x30 celdas de 33px",
            "FPS: 7 (velocidad del juego)",
            "Teclas: flechas para mover"
        ],
        "Persistencia", [
            "scores.json guarda top 10",
            "Formato: [{name, score}, ...]",
            "Se ordena por score descendente",
            "Se actualiza al guardar nuevo score"
        ],
        top=Inches(2))

    # === DIAPOSITIVA 8: RETOS ===
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide8, BG_COLOR)
    add_title(slide8, "Retos del Proyecto")

    add_bullet_list(slide8, [
        "1. Aprender UV → Gestor de dependencias moderno, uv add/remove/sync",
        "2. Configurar Pygame → Inicialización, eventos, loop, pygame.draw",
        "3. Draw.io y Mermaid → Diagramas de flujo y documentación visual",
        "4. Primer proyecto Python → Estructura, módulos, clases, JSON"
    ], Inches(0.5), Inches(1.8), Inches(9), Inches(4))

    add_text(slide8, "Cada reto fue una oportunidad de aprendizaje",
             Inches(0.5), Inches(6), Inches(9), Inches(0.5),
             font_size=Pt(18), color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)

    # === DIAPOSITIVA 9: CONCLUSIÓN ===
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide9, BG_COLOR)
    add_title(slide9, "Conclusión")

    add_bullet_list(slide9, [
        "Estructura clara y organizada del proyecto",
        "Máquina de estados para controlar el flujo",
        "Persistencia de datos con JSON",
        "Documentación visual con diagramas",
        "Herramientas modernas: uv, Pygame"
    ], Inches(0.5), Inches(1.8), Inches(9), Inches(3))

    add_text(slide9, "¡Gracias!", Inches(0.5), Inches(5.5), Inches(9), Inches(1),
             font_size=Pt(40), color=TITLE_COLOR, bold=True, alignment=PP_ALIGN.CENTER)

    # Guardar
    prs.save(OUTPUT_FILE)
    print(f"✓ Presentación guardada: {OUTPUT_FILE}")


if __name__ == "__main__":
    print("=== Generando presentación: Puntos Clave ===\n")
    create_presentation()
